from pptx import Presentation
import json

flag = []
real_key = "real_key/"
slide_name =  "START.pptx"
slide_page = 1

while (True):
    try:
        prs = Presentation(real_key + slide_name)
        sp = prs.slides[slide_page - 1].shapes.title.text.split(", ")
        flag.append(sp[0])
        slide_name = sp[1]
        slide_page = int(sp[2])
    except:
        break

for item in flag:
    print(item, end='')